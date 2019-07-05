from pptx import Presentation
# import glob

def convert_ppt2txt(DocPath, UserDoc):
	prs = Presentation(DocPath)
	page = ""
	for slide in prs.slides:
		for shape in slide.shapes:
			if hasattr(shape,"text"):
				page+=shape.text
				# print(page)

	File = open(UserDoc,"w+", encoding='utf8', errors='ignore')
	File.write(page)
	File.close()


# convert_ppt2txt("C:/Users/rohaa/Desktop/angularFlask/Pdfs/DataScience.pptx","Test_document.txt")
