# -*- coding: utf-8 -*-
"""
Created on Mon May 27 03:06:43 2019

@author: rohaa
"""

from pptx import Presentation

def readPptX(DocPath, UserDoc):
        prs = Presentation(DocPath)
        for slide in prs.slides:
                for shape in slide.shapes:
                        print(shape.text)
                        print("\n\n HI \n\n")

        File_object = open(UserDoc,"w+", encoding='utf8', errors='ignore')
        File_object.write(text)
        File_object.close()
        